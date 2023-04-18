from create_classes import Assignment, Grade, Token, Classroom
from discord.ext import commands
from http import client
from pptx import Presentation
from PyPDF2 import PdfReader
from typing import Optional, List
from supabase import create_client, Client

import api
import asyncio
import create_quiz, create_assignment, create_discussion
import datetime
import discord
import docx
import io
import json
import openai
import os
import PyPDF2
import requests
import secrets
import string
import utils

if os.path.exists(os.getcwd() + "/config.json"):
    with open("./config.json") as f:
        configData = json.load(f)
else:
    configTemp = {"DiscordToken": "", "Prefix": "!", "SupaUrl": "", "SupaKey": ""}
    with open(os.getcwd() + "/config.json", "w+") as f:
        json.dump(configTemp, f)


def run_discord_bot():
    DISCORD_TOKEN = configData["DiscordToken"]
    PREFIX = configData["Prefix"]
    SB_URL = configData["SupaUrl"]
    SB_KEY = configData["SupaKey"]
    GPT_KEY = configData["GPTKey"]

    supabase: Client = create_client(SB_URL, SB_KEY)

    openai.api_key = GPT_KEY

    bot = commands.Bot(command_prefix='!', intents=discord.Intents.all())

    @bot.event
    async def on_ready():
        print(f'{bot.user} is now running!')
        bot.add_view(create_quiz.StartQuiz())

    @bot.event
    async def on_guild_remove(guild: discord.Guild):
        await asyncio.sleep(7 * 24 * 60 * 60)
        current_servers = [guild.id for guild in client.guilds]
        classrooms = api.get_classrooms()
        for classroom in classrooms:
            if classroom.serverId not in current_servers:
                await api.remove_classroom(guild.id)

    @bot.event
    async def on_guild_join(guild: discord.Guild):

        text_channel = guild.text_channels[0]
        message = await text_channel.send(
            f"Warning: {guild.owner.mention} This bot will delete all text/voice channels and roles. Are you sure you "
            "want to add it? React with ✅ to confirm or ❌ to cancel.")
        await message.add_reaction('✅')
        await message.add_reaction('❌')

        def check(reaction, user):
            return user == guild.owner and reaction.message.id == message.id and str(reaction.emoji) in ['✅', '❌']

        while True:
            reaction, user = await bot.wait_for('reaction_add', check=check)

            # Assign role to user who reacted
            if reaction.emoji == "✅":
                break
            if reaction.emoji == "❌":
                return await guild.leave()

        for category in guild.categories:
            await category.delete()

        for channel in guild.channels:
            await channel.delete()

        # Get the bot's member object
        bot_member = guild.get_member(bot.user.id)

        # Get the highest role of the bot
        bot_highest_role = bot_member.top_role
        all_perms = discord.Permissions.all()
        educator_role = await guild.create_role(name="Educator", color=discord.Color(0xffff00), permissions=all_perms,
                                                mentionable=True, hoist=True)
        await educator_role.edit(position=bot_highest_role.position - 1)
        grading_perms = discord.Permissions.all_channel()
        assistant_role = await guild.create_role(name="Assistant", color=discord.Color(0xff8800),
                                                 permissions=grading_perms, mentionable=True, hoist=True)
        student_perms = discord.Permissions.none()
        student_perms.update(
            add_reactions=True, stream=True, read_messages=True, view_channel=True,
            send_messages=True, embed_links=True, attach_files=True, read_message_history=True,
            connect=True, speak=True, use_voice_activation=True, change_nickname=True, use_application_commands=True,
            create_public_threads=True, send_messages_in_threads=True, use_embedded_activites=True
        )
        student_role = await guild.create_role(name="Student", color=discord.Color(0x8affe9), permissions=student_perms,
                                               mentionable=True, hoist=True)

        everyone_perms = discord.Permissions.none()
        everyone_perms.update(
            read_message_history=True, read_messages=True
        )
        everyone = guild.default_role
        await everyone.edit(permissions=everyone_perms)

        # Create server categories
        upcoming = await guild.create_category("Upcoming")
        lounge = await guild.create_category("Lounge")
        await lounge.set_permissions(student_role, view_channel=False)
        await lounge.set_permissions(everyone, view_channel=False)
        general = await guild.create_category("General")
        await guild.create_category("Assignments")
        await guild.create_category("Quizzes")
        await guild.create_category("Discussions")
        submissions = await guild.create_category("Submissions")
        await submissions.set_permissions(student_role, view_channel=False)
        await submissions.set_permissions(everyone, view_channel=False)
        grades = await guild.create_category("Grades")
        await grades.set_permissions(student_role, view_channel=False)
        await grades.set_permissions(everyone, view_channel=False)
        questions = await guild.create_category("Questions")

        announcement_channel = await guild.create_text_channel("Announcements", category=general)
        await announcement_channel.set_permissions(student_role, send_messages=False)
        await guild.create_text_channel("General", category=general)
        attendance_channel = await guild.create_text_channel("Attendance", category=general)
        await attendance_channel.set_permissions(everyone, view_channel=False)
        await attendance_channel.set_permissions(student_role, view_channel=True, send_messages=False,
                                                 add_reactions=False)
        syllabus_channel = await guild.create_text_channel("Syllabus", category=general)
        await syllabus_channel.set_permissions(student_role, send_messages=False)
        await guild.create_voice_channel("Lecture", category=general)
        await guild.create_voice_channel("Chill", category=general)

        await guild.create_text_channel("Educators-Assistants", category=lounge)
        await guild.create_text_channel("Terminal", category=lounge)
        await guild.create_voice_channel("Educators-Assistants", category=lounge)

        await guild.create_text_channel("Public", category=questions)

        classroom = Classroom(attendance=0, serverId=guild.id, serverName=guild.name)
        await api.create_classroom(classroom)

        # gives discord owner the Educator role
        await guild.owner.add_roles(educator_role)

    # Gives new users the Student role
    @bot.event
    async def on_member_join(member):
        role = discord.utils.get(member.guild.roles, name="Student")
        await member.add_roles(role)
        discordNickname = member.display_name
        discordId = member.id
        await utils.add_member_to_table(guild_id=member.guild.id, role="Student", nick=discordNickname,
                                        discord_id=discordId)

    @bot.event
    async def on_member_remove(member: discord.Member):
        if not member.bot:
            res = await api.get_classroom_id(member.guild.id)
            classroom_id = res.content['id']
            res = await api.get_user_id(member.id)
            user_id = res.content['id']

            # TODO: implement user delete
            supabase.table("Classroom_User").delete().match({"classroomId": classroom_id, "userId": user_id}).execute()

    @bot.event
    async def on_member_update(before: discord.Member, after: discord.Member):
        # Update member nickname in database
        if before.nick != after.nick:
            await api.update_classroom_user_name(after.nick, after.id)

        # Update member role in database
        if before.roles != after.roles:
            role_names = ['Student', 'Educator', 'Assistant']
            guild = before.guild
            roles = {role.name: role for role in guild.roles if role.name in role_names}
            before_member_roles = [role for role in before.roles if role.name in role_names]
            after_member_roles = [role for role in after.roles if role.name in role_names]
            role_count = len(after_member_roles)
            res = await api.get_user_id(after.id)
            user_id = res['id']
            server_id = after.guild.id
            res = await api.get_classroom_id(server_id)
            classroom_id = res['id']
            if role_count == 0:
                await after.add_roles(roles['Student'])
                new_role = 'Student'
                if not any(role.name == "Student" for role in before.roles):
                    await api.update_classroom_user_role(user_id, classroom_id, new_role)
            elif role_count > 1:
                new_roles = [x for x in after_member_roles if x not in before_member_roles]
                new_role = new_roles[0].name
                latest_role = new_roles[0]
                roles_to_remove = [role for role in after_member_roles if role != latest_role]

                for role in roles_to_remove:
                    await after.remove_roles(role)

                await api.update_classroom_user_role(user_id, classroom_id, new_role)  # TODO: API

    @bot.slash_command(name='syllabus',
                       description='```/syllabus [.pdf file]``` - Updates the syllabus page with the linked pdf file')
    async def syllabus(ctx: discord.ApplicationContext, file: discord.Attachment):
        if ctx.author.id != ctx.guild.owner_id:
            await ctx.send("Error, enter '!help' for  more information.")
            return
        owner = ctx.guild.owner
        # restrict syllabus channel but allow for visibility.
        overwrites = {
            ctx.guild.default_role: discord.PermissionOverwrite(read_messages=True, send_messages=False),
            owner: discord.PermissionOverwrite(read_messages=True, send_messages=True)
        }

        await ctx.respond("Success", delete_after=0)

        if file is not None:
            pdf = file
            if pdf.filename.endswith(".pdf"):
                file_data = await pdf.read()
                # Check whether "syllabus" channel exists;
                channel = discord.utils.get(ctx.guild.channels, name="syllabus")
                if channel is None:
                    channel = await ctx.guild.create_text_channel("syllabus", overwrites=overwrites)

                await channel.send("**```diff\n+ Class Syllabus```**")
                pdf_reader = PdfReader(io.BytesIO(file_data))
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                text_bytes = text.encode('utf-8')
                max_char = 1500
                # append and exclude unicodeerror
                chunks = [text_bytes[i:i + max_char].decode('utf-8', errors='ignore') for i in
                          range(0, len(text_bytes), max_char)]

                # send text representation of syllabus to syllabus channel
                for chunk in chunks:
                    await channel.send("```" + chunk + "```")
                await channel.send(
                    "**```diff\n- Please note, the text representation of the syllabus may not be completely accurate. For a more precise version, download syllabus provided below.```**")

                # Send PDF "syllabus" channel
                await channel.send(file=discord.File(io.BytesIO(file_data), filename=pdf.filename))
                await ctx.send("Syllabus uploaded successfully.")

            else:
                await ctx.send("Invalid file format, only PDF files are accepted.")

    @bot.slash_command(name='delete', description='```/delete ``` - Deletes a quiz, assignment, or discussion')
    async def delete(ctx: discord.ApplicationContext):
        user_roles = [role.name for role in ctx.author.roles]

        if "Educator" not in user_roles:
            await ctx.respond("You do not have permission to use this command. Only an educator can use /delete",
                              delete_after=10)
        else:
            channel = ctx.channel
            embed = discord.Embed(title="Delete this channel?", color=0x00FF00)
            embed.add_field(name="Are you sure you want to delete this channel?", value="✅ Yes\n❌ No")
            interaction: discord.Interaction = await ctx.respond(embed=embed)
            message: discord.Message = await interaction.original_response()
            await message.add_reaction("✅")
            await message.add_reaction("❌")

            def check(reaction, user):
                return user == ctx.author and reaction.message.id == message.id and str(reaction.emoji) in ["✅", "❌"]

            while True:
                reaction, user = await bot.wait_for('reaction_add', check=check)

                if reaction.emoji == "✅":
                    await channel.delete()
                elif reaction.emoji == "❌":
                    await message.delete()

    @bot.slash_command(name='poll',
                       description='```/poll [topic] [option1] [option2] ... [option8]``` - Creates a poll for users (8 max options)')
    async def poll(ctx: discord.ApplicationContext, topic: str, option1: str, option2: str, option3: str = None,
                   option4: str = None, option5: str = None, option6: str = None, option7: str = None,
                   option8: str = None):

        await ctx.defer()

        options = [option1, option2]
        if option3:
            options.append(option3)
        if option4:
            options.append(option4)
        if option5:
            options.append(option5)
        if option6:
            options.append(option6)
        if option7:
            options.append(option7)
        if option8:
            options.append(option8)

        num_reactions_for_option_idx = [0] * len(options)

        # box area created for reactions
        boxes = []
        for i in range(10):
            boxes.append("⬛")
        boxesStr = ''.join(boxes)

        # Create the poll embed
        embed = discord.Embed(title=topic, description=' '.join(
            [f'{chr(0x1f1e6 + i)} {option}\n `{boxesStr}` {0} (0%)\n' for i, option in enumerate(options)]))

        await ctx.respond("Poll Created")

        # Send the poll message and add reactions
        message = await ctx.send(embed=embed)
        for i in range(len(options)):
            await message.add_reaction(chr(0x1f1e6 + i))

        num_total_reactions = [0]

        async def update_poll_data(change_to_reactions, num_total_reactions, reaction):
            num_total_reactions[0] += change_to_reactions
            # print(f"{num_total_reactions[0]} users reacted to the message")

            new_description_for_option_idx = [None] * len(options)
            for i in range(len(options)):
                # If the reaction emoji corresponds to options[i],
                # update num_reactions_for_option_idx[i]
                if reaction.emoji == chr(127462 + i):
                    num_reactions_for_option_idx[i] += change_to_reactions

                # print(f"option{i + 1} reaction total: {num_reactions_for_option_idx[i]}")
                reactionPercentage = round((num_reactions_for_option_idx[i] / num_total_reactions[0]) * 100) if \
                    num_total_reactions[0] > 0 else 0
                # print(f"option{i + 1}%: {reactionPercentage}")

                boxesForOption = boxes.copy()
                for j in range(round(reactionPercentage / 10)):
                    if i < 7:
                        boxesForOption[j] = f"{chr(0x1F7E5 + i)}"
                    else:
                        boxesForOption[j] = f"{chr(0x2B1C)}"
                boxesForOptionStr = ''.join(boxesForOption)

                new_description_for_option_idx[
                    i] = f'{chr(0x1f1e6 + i)} {options[i]}\n `{boxesForOptionStr}` {num_reactions_for_option_idx[i]} ({reactionPercentage}%)\n'

            embed.description = ' '.join(new_description_for_option_idx)
            await message.edit(embed=embed)

        # check that the reaction is from a user, not a bot
        def check(reaction, user):
            return user != bot.user and reaction.message.id == message.id and str(reaction.emoji) in [
                chr(127462 + i) for i in range(len(options))]

        async def on_reaction_add(reaction, user):
            if check(reaction, user):
                await update_poll_data(1, num_total_reactions, reaction)

        async def on_reaction_remove(reaction, user):
            await update_poll_data(-1, num_total_reactions, reaction)

        # bot waits for reaction event from user
        bot.add_listener(on_reaction_add, 'on_reaction_add')
        bot.add_listener(on_reaction_remove, 'on_reaction_remove')

    class AnonPoll(discord.ui.View):  # class is created for the discord view of the poll with buttons
        def __init__(self, options: List[str], ctx: discord.ApplicationContext, message: discord.Message,
                     user: discord.User):
            super().__init__()
            self.creator = user
            self.options = options
            self.buttons = [None] * len(options)
            self.num_votes_for_option_idx = [0] * len(options)  # number of votes for each option
            self.num_total_votes = 0
            self.votes = {}

            def makeOptionCallback(idx):  # callback function for when user presses button to select an option
                async def optionCallback(interaction: discord.Interaction):
                    if interaction.user.id in self.votes[idx]:
                        self.votes[idx].remove(interaction.user.id)
                        self.num_total_votes -= 1
                        self.num_votes_for_option_idx[idx] -= 1
                        await interaction.response.send_message(f'Removed vote for {self.options[idx]}', ephemeral=True)
                    else:
                        self.votes[idx].append(interaction.user.id)
                        self.num_votes_for_option_idx[idx] += 1
                        self.num_total_votes += 1
                        await interaction.response.send_message(f'You voted for {self.options[idx]}', ephemeral=True)

                return optionCallback

            for i in range(len(self.options)):  # each option is created to be a button
                button = discord.ui.Button(label=self.options[i], style=discord.ButtonStyle.secondary,
                                           emoji=chr(0x1f1e6 + i))
                self.buttons[i] = button
                self.votes[i] = []
                button.callback = makeOptionCallback(i)
                self.add_item(button)

            async def endPollCallback(
                    interaction: discord.Interaction):  # callback function to display results when poll ends
                if interaction.user != self.creator:
                    return await interaction.response.send_message("Only the creator of the poll can end it",
                                                                   ephemeral=True)

                description_for_option_idx = [None] * len(options)

                for i in range(len(options)):  # disply the poll results for each option
                    boxesForOption = ["⬛"] * 10

                    reactionPercentage = round((self.num_votes_for_option_idx[
                                                    i] / self.num_total_votes) * 100) if self.num_total_votes > 0 else 0

                    for j in range(round(reactionPercentage / 10)):
                        if i < 7:
                            boxesForOption[j] = f"{chr(0x1F7E5 + i)}"
                        else:
                            boxesForOption[j] = f"{chr(0x2B1C)}"

                    boxesForOptionStr = ''.join(boxesForOption)

                    description_for_option_idx[
                        i] = f'{chr(0x1f1e6 + i)} {options[i]}\n `{boxesForOptionStr}` {self.num_votes_for_option_idx[i]} ({reactionPercentage}%)\n'

                embed = message.embeds[0]  # gets the poll message embed
                embed.description = ' '.join(description_for_option_idx)
                await message.edit(embed=embed)

                endPollButton.disabled = True
                for button in self.buttons:
                    button.disabled = True

                await interaction.response.edit_message(view=self)

            # get the roles for each user
            user_roles = [role.name for role in ctx.author.roles]
            endPollButton = discord.ui.Button(label="End Poll", style=discord.ButtonStyle.danger)
            endPollButton.callback = endPollCallback
            self.add_item(endPollButton)

    @bot.slash_command(name='anonpoll',
                       description='```/anon poll [topic] [option1] [option2] ... ``` - Creates a poll for users (8 max options)')
    async def anon_poll(ctx: discord.ApplicationContext, topic: str, option1: str, option2: str, option3: str = None,
                        option4: str = None, option5: str = None, option6: str = None, option7: str = None,
                        option8: str = None):

        await ctx.defer()

        options = [option1, option2]
        if option3:
            options.append(option3)
        if option4:
            options.append(option4)
        if option5:
            options.append(option5)
        if option6:
            options.append(option6)
        if option7:
            options.append(option7)
        if option8:
            options.append(option8)

        # create the poll embed to show user options
        embed = discord.Embed(title=topic, description=' '.join(
            [f'{chr(0x1f1e6 + i)} {option}\n\n' for i, option in enumerate(options)]))

        await ctx.respond("Poll Created")

        # Send the poll message and add buttons to the view
        message = await ctx.send(embed=embed)
        await message.edit(view=AnonPoll(options, ctx, message, ctx.user))

    @bot.slash_command(name='attendance',
                       description='```/attendance``` - Used by students to check their attendance')
    async def attendance(ctx: discord.ApplicationContext):
        user_roles = [role.name for role in ctx.author.roles]
        guild_id = ctx.guild_id
        if "Student" not in user_roles:
            return await ctx.respond(
                "Only Students can use /attendance\nIf you're trying to take attendance use `/lecture attendance`")
        res = supabase.table('User').select('id').eq('discordId', ctx.author.id).execute()
        user_id = res.data[0]['id']
        res = await api.get_classroom_id(ctx.guild_id)
        classroom_id = res.content['id']
        res = supabase.table('Classroom_User').select('attendance').match(
            {'userId': user_id, 'classroomId': classroom_id}).execute()
        user_attendance = res.data[0]['attendance']
        res = supabase.table('Classroom').select('attendance').eq('serverId', ctx.guild_id).execute()
        classroom_attendance = res.data[0]['attendance']
        response = f"Your attendance count is {user_attendance} / {classroom_attendance}."
        await ctx.author.send(response)
        await ctx.respond("Check DMs for your attendance", delete_after=2)

    @bot.slash_command(name='ta',
                       description='```/ta [user]``` - Gives/Removes the user the Assistant role')
    async def ta(ctx: discord.ApplicationContext, user: discord.Member):
        edu_role = discord.utils.get(ctx.guild.roles, name="Educator")
        if edu_role in ctx.author.roles:
            role = discord.utils.get(ctx.guild.roles, name="Assistant")
            if role in user.roles:
                await user.remove_roles(role)
                await ctx.respond(f"Removed Assistant role from {user.mention}")
            else:
                await user.add_roles(role)
                await ctx.respond(f"{user.mention} has been given the Assistant role")
        else:
            await ctx.respond("You need the Educator role to use this command")

    @bot.slash_command(name='edu',
                       description='```/edu [user]``` - Gives/Removes the user the Educator role')
    async def edu(ctx: discord.ApplicationContext, user: discord.Member):
        edu_role = discord.utils.get(ctx.guild.roles, name="Educator")
        if edu_role in ctx.author.roles:
            if user == ctx.guild.owner:
                return await ctx.respond("You cannot remove the Educator role from the owner")
            if edu_role in user.roles:
                await user.remove_roles(edu_role)
                await ctx.respond(f"Removed Educator role from {user.mention}")
            else:
                await user.add_roles(edu_role)
                await ctx.respond(f"{user.mention} has been given the Educator role")
        else:
            await ctx.respond("You need the Educator role to use this command")

    @bot.slash_command(name='private',
                       description='```/private [question]``` - Creates a private question between Student and Educator/TA')
    async def private(ctx: discord.ApplicationContext, question: str):
        student_role = discord.utils.get(ctx.guild.roles, name="Student")
        if student_role in ctx.author.roles:
            user_mention = ctx.author.mention
            student_id = str(ctx.author.id)
            q = discord.utils.get(ctx.guild.categories, name="Questions")
            if q is None:
                q = await ctx.guild.create_category("Questions")

            channels = q.channels

            for channel in channels:
                if channel.name == f"private-{student_id}":
                    await ctx.respond(f"{user_mention} You already have a private question open.", delete_after=3)
                    return

            private_role = await ctx.guild.create_role(name=f"{student_id}")

            await ctx.author.add_roles(private_role)

            all_roles = ctx.guild.roles

            ow = {}

            for role in all_roles:
                ow[role] = discord.PermissionOverwrite(read_messages=False)

            overwrites = {
                ctx.guild.default_role: discord.PermissionOverwrite(read_messages=False),
                private_role: discord.PermissionOverwrite(read_messages=True),
                discord.utils.get(ctx.guild.roles, name="Educator"): discord.PermissionOverwrite(read_messages=True),
                discord.utils.get(ctx.guild.roles, name="Assistant"): discord.PermissionOverwrite(read_messages=True)
            }

            ow.update(overwrites)

            private_channel = await ctx.guild.create_text_channel(f"Private-{student_id}", category=q, overwrites=ow)

            await ctx.respond("Private question created", delete_after=3)

            await private_channel.send(f"{user_mention} asked: {question}")
        else:
            await ctx.respond("Only Students can ask private questions")

    @bot.slash_command(name='close', description='```/close``` - Deletes a private question channel')
    async def close(ctx: discord.ApplicationContext):
        channel_name = ctx.channel.name
        category = ctx.channel.category.name
        if category == "Questions" and 'private' in channel_name:
            embed = discord.Embed(title="Close this question?", color=0x00FF00)
            embed.add_field(name="Are you sure you want to close this question?", value="✅ Yes\n❌ No")
            interaction: discord.Interaction = await ctx.respond(embed=embed)
            message: discord.Message = await interaction.original_response()
            await message.add_reaction("✅")
            await message.add_reaction("❌")

            def check(reaction, user):
                return user == ctx.author and reaction.message.id == message.id and str(reaction.emoji) in ["✅", "❌"]

            while True:
                reaction, user = await bot.wait_for('reaction_add', check=check)

                # Assign role to user who reacted
                if reaction.emoji == "✅":
                    student_id = ctx.channel.name.split("-")[1]
                    created_role = discord.utils.get(ctx.guild.roles, name=f"{student_id}")
                    await created_role.delete()
                    await ctx.channel.delete(reason=f"Closed by {ctx.author.display_name}")
                if reaction.emoji == "❌":
                    await message.delete()
        else:
            await ctx.respond("You cannot close this channel", delete_after=3)

    @bot.slash_command(name='help', description='```/help``` - Sends command information to the user')
    async def help(ctx: discord.ApplicationContext, command_name: str = None):
        if command_name:
            command = bot.get_application_command(command_name)
            if not command:
                await ctx.respond(f"There exists no command named '{command_name}'.")
                return
            await ctx.respond(f"**{command_name}:**\n{command.description}")
        else:
            everyone_commands = ['tutor', 'help', 'poll', 'anonpoll', 'close']
            student_commands = ['submit',  'attendance', 'private', 'upload_file'] + everyone_commands
            educator_commands = ['create', 'grade', 'lecture', 'delete', 'edu', 'ta', 'syllabus', 'upcoming', 'edit'] + everyone_commands
            assistant_commands = ['grade', 'edit', 'lecture', 'delete', 'upcoming'] + everyone_commands
            user_roles = [role.name for role in ctx.user.roles]
            my_commands = []
            if 'Student' in user_roles:
                my_commands = student_commands
            elif 'Educator' in user_roles:
                my_commands = educator_commands
            elif 'Assistant' in user_roles:
                my_commands = assistant_commands
            else:
                my_commands = everyone_commands
            message = f"**Available Commands for {user_roles[1]} role:**\n\n"
            for command in bot.application_commands:
                if command.name not in my_commands:
                    continue

                if command.name == "create" and 'Educator' in user_roles:
                    for create in command.walk_commands():
                        message += f"{create.description}\n\n"

                elif command.name == 'lecture' and ('Educator' in user_roles or 'Assistant' in user_roles):
                    for lecture in command.walk_commands():
                        message += f"{lecture.description}\n\n"

                elif command.name == 'grade' and ('Educator' in user_roles or 'Assistant' in user_roles):
                    for grade in command.walk_commands():
                        message += f"{grade.description}\n\n"
                elif command.name == 'tutor':
                    for tutor in command.walk_commands():
                        message += f"{tutor.description}\n\n"
                else:
                    message += f"{command.description}\n\n"

            await ctx.author.send(message)
            await ctx.respond("Check Direct Messages for available commands")

    create = bot.create_group("create", "create school work")

    @create.command(name='quiz',
                    description='```/create quiz``` - Creates a Quiz for students to take')
    async def quiz(ctx: discord.ApplicationContext):
        edu_role = discord.utils.get(ctx.guild.roles, name="Educator")
        if edu_role in ctx.author.roles:
            modal = create_quiz.create_quiz(bot=bot)
            await ctx.send_modal(modal)
        else:
            await ctx.respond("You need to be an Educator to create quizzes", delete_after=3)

    @create.command(name='discussion',
                    description='```/create discussion``` - Creates a new text channel with a prompt for discussion')
    async def discussion(ctx: discord.ApplicationContext):
        edu_role = discord.utils.get(ctx.guild.roles, name="Educator")
        if edu_role in ctx.author.roles:
            modal = create_discussion.create_discussion(bot=bot)
            await ctx.send_modal(modal)
        else:
            await ctx.respond("You need to be an Educator to create discussions", delete_after=3)

    @create.command(name='assignment',
                    description='```/create assignment (file)``` - Creates assignments for students')
    async def assignment(ctx: discord.ApplicationContext, file: discord.Attachment = None):
        edu_role = discord.utils.get(ctx.guild.roles, name="Educator")
        if edu_role in ctx.author.roles:
            modal = create_assignment.create_assignment(bot=bot, file=file)
            await ctx.send_modal(modal)
        else:
            await ctx.respond("You need to be an Educator to create assignments", delete_after=3)

    @create.command(name='upload', description='```/create upload [json file]``` - Creates a task from a json file')
    async def upload(ctx: discord.ApplicationContext, file: discord.Attachment):
        edu_role = discord.utils.get(ctx.guild.roles, name="Educator")
        if edu_role in ctx.author.roles:
            contents = await file.read()
            data = json.loads(contents)

            if data['type'] == "Discussion":
                modal = create_discussion.create_discussion(bot=bot, preset=data)
                await ctx.send_modal(modal)
            if data['type'] == "Assignment":
                modal = create_assignment.create_assignment(bot=bot, preset=data)
                await ctx.send_modal(modal)
            if data['type'] == "Quiz":
                modal = create_quiz.create_quiz(bot=bot, preset=data)
                await ctx.send_modal(modal)
        else:
            await ctx.respond("You need to be an Educator to use /create upload", delete_after=3)

    lecture = bot.create_group("lecture", "Commands used during lectures")

    @lecture.command(name="attendance",
                     description='```/lecture attendance [time (minutes)]``` - Creates attendance poll')
    async def attendance(ctx: discord.ApplicationContext, time: float = 5):
        user_roles = [role.name for role in ctx.author.roles]
        guild_id = ctx.guild_id
        if "Educator" not in user_roles:
            return await ctx.respond("Only educators can use /lecture commands")
        attendance_channel = discord.utils.get(ctx.guild.channels, name="attendance")
        await ctx.respond("Now taking Attendance...", delete_after=3)
        date = datetime.datetime.now().strftime("%m - %d - %y %I:%M %p")
        student_role = discord.utils.get(ctx.guild.roles, name="Student")
        embed = discord.Embed(title="Attendance",
                              description=f'{student_role.mention} React to this message to check into today\'s attendance')
        message = await attendance_channel.send(embed=embed)
        await message.add_reaction('✅')

        timeLeft = time * 60
        while timeLeft >= 0:
            embed.title = f"Attendance - {int(timeLeft)}s"
            await asyncio.sleep(1)
            await message.edit(embed=embed)
            timeLeft -= 1
        embed.description = "Attendance CLOSED"
        await asyncio.sleep(1)
        await message.edit(embed=embed)
        attendance_message = await attendance_channel.fetch_message(message.id)
        reactions = attendance_message.reactions
        users = []
        for r in reactions:
            if r.emoji == '✅':
                async for user in r.users():
                    users.append(user)

        attended = []
        for user in users:
            if not user.bot:
                await utils.increment_attendance(user.id, ctx.guild_id)
                if user.nick is not None:
                    attended.append(user.nick)
                else:
                    attended.append(user.name)

        response = f"Attendance for {date}:\n\nAttended:\n" + '\n'.join(attended)

        students = discord.utils.get(ctx.guild.roles, name="Student").members

        absent = []

        for student in students:
            if student not in users:
                if student.nick is not None:
                    absent.append(student.nick)
                else:
                    absent.append(student.name)

        response += "\n\nAbsent:\n" + '\n'.join(absent)
        await ctx.author.send(response)

        # Update the 'total attendance' in the Supabase table
        request = supabase.table('Classroom').select('attendance').eq('serverId', ctx.guild_id).execute()
        classroom_attendance = request.data[0]['attendance']
        _, error = supabase.table('Classroom').update({'attendance': classroom_attendance + 1}).eq('serverId',
                                                                                                   ctx.guild_id).execute()

    @lecture.command(name="mute",
                     description='```/lecture mute``` - Mutes everyone in the voice channel besides the Educator')
    async def mute(ctx: discord.ApplicationContext):
        user_roles = [role.name for role in ctx.author.roles]
        if "Educator" not in user_roles:
            return await ctx.respond("Only educators can use /lecture commands", delete_after=3)
        if ctx.author.voice is None:
            return await ctx.respond("You are not in a voice channel", delete_after=3)
        voice_channel = ctx.author.voice.channel
        members_in_voice = voice_channel.members

        for member in members_in_voice:
            if member != ctx.author:
                await member.edit(mute=True)

        await ctx.respond("Students Muted.\nUse `/lecture unmute` to remove the mute from the members", delete_after=3)

    @lecture.command(name="unmute",
                     description='```/lecture unmute``` - Unmutes everyone in the voice channel')
    async def unmute(ctx: discord.ApplicationContext):
        user_roles = [role.name for role in ctx.author.roles]
        if "Educator" not in user_roles:
            return await ctx.respond("Only educators can use /lecture commands", delete_after=3)
        if ctx.author.voice is None:
            return await ctx.respond("You are not in a voice channel", delete_after=3)
        voice_channel = ctx.author.voice.channel
        members_in_voice = voice_channel.members

        for member in members_in_voice:
            if member != ctx.author:
                await member.edit(mute=False)

        await ctx.respond("Students Unmuted.", delete_after=3)

    @lecture.command(name="breakout",
                     description='```/lecture breakout [number_of_rooms]``` - Creates a number of breakout rooms for students')
    async def breakout(ctx: discord.ApplicationContext, number_of_rooms: int, number_of_minutes: float = 60):
        general_category = discord.utils.get(ctx.guild.categories, name="General")

        rooms = []

        for i in range(number_of_rooms):
            rooms.append(await ctx.guild.create_voice_channel(f"Breakout Room {i + 1}", category=general_category))

        interaction = await ctx.respond(
            f"{number_of_rooms} Breakout rooms created, react with ❌ to remove them when done.")
        message = await interaction.original_response()
        await message.add_reaction('❌')

        def check(reaction, user):
            return user == ctx.author and reaction.message.id == message.id and str(reaction.emoji) in ['❌']

        while True:
            try:
                reaction, user = await bot.wait_for('reaction_add', check=check, timeout=number_of_minutes * 60)
                for room in rooms:
                    await room.delete()
                break
            except asyncio.TimeoutError:
                for room in rooms:
                    await room.delete()
                break

        await message.edit("Breakout rooms deleted.", delete_after=3)

    @bot.slash_command(name='edit',
                       description='```/edit``` - Used by Educators to edit quizzes, assignments, and discussions')
    async def edit(ctx: discord.ApplicationContext):
        edu_role = discord.utils.get(ctx.guild.roles, name="Educator")
        if edu_role in ctx.author.roles:

            if ctx.channel.category.name == "Assignments":
                assignment_dict = {}
                first_message = await ctx.channel.history(oldest_first=True, limit=1).next()
                assignment_dict['type'] = first_message.embeds[0].title
                assignment_dict['title'] = first_message.embeds[0].fields[0].value
                assignment_dict['details'] = first_message.embeds[0].fields[1].value
                assignment_dict['points'] = first_message.embeds[0].fields[2].value
                assignment_dict['start'] = first_message.embeds[0].fields[3].value
                assignment_dict['due'] = first_message.embeds[0].fields[4].value
                assignment_dict['channelId'] = ctx.channel_id
                modal = create_quiz.EditModal(assignment_dict, first_message, title="Editing Assignment")
                await ctx.send_modal(modal)
            elif ctx.channel.category.name == "Quizzes":
                quiz_dict = {}
                first_message = await ctx.channel.history(oldest_first=True, limit=1).next()
                quiz_dict['type'] = first_message.embeds[0].title
                quiz_dict['title'] = first_message.embeds[0].fields[0].value
                quiz_dict['time'] = first_message.embeds[0].fields[1].value
                quiz_dict['points'] = first_message.embeds[0].fields[2].value
                quiz_dict['start'] = first_message.embeds[0].fields[3].value
                quiz_dict['due'] = first_message.embeds[0].fields[4].value
                quiz_dict['channelId'] = ctx.channel_id
                modal = create_quiz.EditModal(quiz_dict, first_message, title="Editing Quiz")
                await ctx.send_modal(modal)
            elif ctx.channel.category.name == "Discussions":
                discussion_dict = {}
                first_message = await ctx.channel.history(oldest_first=True, limit=1).next()
                discussion_dict['type'] = first_message.embeds[0].title
                discussion_dict['title'] = first_message.embeds[0].fields[0].value
                discussion_dict['details'] = first_message.embeds[0].fields[1].value
                discussion_dict['points'] = first_message.embeds[0].fields[2].value
                discussion_dict['start'] = first_message.embeds[0].fields[3].value
                discussion_dict['due'] = first_message.embeds[0].fields[4].value
                discussion_dict['channelId'] = ctx.channel_id
                modal = create_quiz.EditModal(discussion_dict, first_message, title="Editing Discussion")
                await ctx.send_modal(modal)
            else:
                await ctx.respond("Can't use /edit here", delete_after=3)
        else:
            await ctx.respond("You need to be an Educator to use /edit", delete_after=3)

    @bot.slash_command(name='upload_file', description='```/upload file``` - User can follow link to upload file')
    async def upload_file(ctx):

        unique_id = ''.join(secrets.choice(string.ascii_letters + string.digits) for _ in range(8, 51))

        discord_id = ctx.author.id
        user = await api.get_user_id(discord_id)
        userId = user.id

        submit_dict = {'userId': userId, 'unique_id': unique_id}

        new_submit = Token(userId=submit_dict['userId'], unique_id=submit_dict['unique_id'])

        await api.update_token(new_submit)

        await ctx.respond(content=f'https://discord-classroom-file-uploads.herokuapp.com/?token={unique_id}',
                          ephemeral=True)

    tutor = bot.create_group("tutor", "AI tutor for students")

    @tutor.command(name='quiz', description='```/tutor quiz [number_of_questions] [subject] [grade]``` - creates quiz')
    async def quiz(ctx: discord.ApplicationContext, number_of_questions: int, subject: str, grade: str):

        if number_of_questions > 5:
            return await ctx.respond("Max number of questions is 5")

        try:
            ordinal_grade = utils.get_ordinal_number(grade)
            # print(f"Ordinal number: {ordinal_grade}")
        except ValueError as e:
            print(e)

        messages = [
            {"role": "system",
             "content": "You are TutorGPT, a friendly and helpful AI that assists students with learning and understanding their school work. You are "
                        "responsible for creating quizzes for students with a given ordinal grade, subject, and number of questions. You should respond in less than "
                        "150 words with the questions. Then I will tell you to generate answers for the questions to then compare with the user's answers. Do not "
                        "respond with intro messages like, just respond with the questions or answers themselves"}
        ]

        await ctx.defer(ephemeral=True)

        messages.append(
            {"role": "user",
             "content": f"Give me a {ordinal_grade} grade quiz on {subject} with {number_of_questions} questions. Do NOT show me the answers."}
        )

        # print(messages)
        chat = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=messages, max_tokens=150)
        reply = chat.choices[0].message.content

        messages.append(
            {"role": 'assistant', 'content': reply}
        )

        print("Quiz:", reply)

        user = ctx.author
        await ctx.respond("Check DMs", delete_after=3)
        await user.send(f"TutorGPT:\n {reply}")

        messages.append(
            {"role": 'user',
             'content': "Now give your answers to the questions above, I won't be able to see your answers so don't worry about cheating."}
        )

        chat = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=messages, max_tokens=150)
        answers = chat.choices[0].message.content

        print("Answers:", answers)

        def check(message):
            return message.author == user and message.channel == user.dm_channel

        response = await bot.wait_for('message', check=check)

        await user.send(f"TutorGPT: Here are the correct answers\n {answers}")

    @tutor.command(name='flashcards',
                   description='```/tutor flashcards [topic]``` creates a list of flashcards on a topic or from notes')
    async def flashcards(ctx: discord.ApplicationContext, topic: str):
        messages = [
            {"role": "system",
             "content": "You are TutorGPT, a friendly and helpful AI that assists students with learning and understanding their school work."}
        ]

        await ctx.defer(ephemeral=True)

        messages.append(
            {"role": "user",
             "content": f"In less than 150 words, can you create 3 flash cards for the following topic or information with the following format:\n"
                        f"Question: INSERT QUESTION HERE\n||Answer: INSERT ANSWER HERE|| for each flashcard. Please do not forget the || before"
                        f"and after the answer, it is very important. If can't finish a flashcards because of the word limit, do not create the"
                        f"flashcard. Here is the topic or information, to"
                        f"save words write the answer as short as possible: {topic}"}
        )

        while True:
            chat = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=messages, max_tokens=200)
            reply = chat.choices[0].message.content
            user = ctx.author
            await ctx.respond("Check DMs", delete_after=3)
            await user.send(f"TutorGPT: {reply}")
            await user.send("If you want more flashcards type 'More'")

            def check(message: discord.Message):
                return message.author == ctx.author and message.channel == ctx.author.dm_channel and message.content == "More"

            try:
                response = await bot.wait_for('message', check=check, timeout=60)
            except asyncio.TimeoutError:
                return await user.send('TutorGPT: 😴')

            messages.append(
                {'role': 'user',
                 'content': 'Can you give me 3 more flashcards that are different then the ones you already gave please'}
            )

    @tutor.command(name='study',
                   description='```/tutor study [file]``` creates a study guide for you from the file you uploaded')
    async def study(ctx: discord.ApplicationContext, file: discord.Attachment):
        messages = [
            {"role": "system",
             "content": "You are TutorGPT, a friendly and helpful AI that assists students with learning and understanding their school work."}
        ]
        await ctx.defer(ephemeral=True)
        supported_files = ['docx', 'pdf', 'pptx']
        file_extension = file.filename.split('.')[-1]
        if file_extension not in supported_files:
            return await ctx.respond("File type not supported.\nSupported File types are: .docx, .pdf, and .pptx")
        information = ""
        if file_extension in ['pptx']:
            binary_data = await file.read()
            ppt = Presentation(io.BytesIO(binary_data))

            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        information += shape.text_frame.text
                        information += "\n"

        elif file_extension == 'docx':
            binary_data = await file.read()
            doc = docx.Document(io.BytesIO(binary_data))

            for paragraph in doc.paragraphs:
                information += paragraph.text
                information += "\n"

        elif file_extension == 'pdf':
            binary_data = await file.read()
            pdf = PyPDF2.PdfReader(io.BytesIO(binary_data))

            for page_number in range(len(pdf.pages)):
                text = pdf.pages[page_number].extract_text()
                information += text
                information += "\n"
        if information == "":
            return await ctx.respond("Could not read contents of file")

        messages.append(
            {"role": "user",
             "content": f"In less than 325 words, can you create a study guide for me as a student with the information I "
                        f"provide you. Ignore information that does not relate to the subject. Here is the information: {information}"}
        )
        try:
            chat = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=messages, max_tokens=330)
        except:
            return await ctx.respond("File contents too large, max words is 3000")

        reply = chat.choices[0].message.content
        user = ctx.author
        await ctx.respond("Check DMs", delete_after=3)
        await user.send(f"TutorGPT: \n{reply}")

    async def get_dates_assignments(server_id: str, due_date: str):
        response = supabase.table("Assignment").select('*').lte('dueDate', due_date).execute()
        all_assignments = response.data
        classroom_assignments = []
        for assignment in all_assignments:
            channel = bot.get_channel(int(assignment['channelId']))
            if channel is not None:
                if str(channel.guild.id) == str(server_id):
                    classroom_assignments.append(assignment)
        return classroom_assignments

    async def get_dates_quiz(server_id: str, due_date: str):
        response = supabase.table("Quiz").select('*').lte('dueDate', due_date).execute()
        all_quizzes = response.data
        classroom_quizzes = []
        for quiz in all_quizzes:
            channel = bot.get_channel(int(quiz['channelId']))
            if channel is not None:
                if str(channel.guild.id) == str(server_id):
                    classroom_quizzes.append(quiz)
        return classroom_quizzes

    async def get_dates_discussion(server_id: str, due_date: str):
        response = supabase.table("Discussion").select('*').lte('dueDate', due_date).execute()
        all_discussions = response.data
        classroom_discussions = []
        for disc in all_discussions:
            channel = bot.get_channel(int(disc['channelId']))
            if channel is not None:
                if str(channel.guild.id) == str(server_id):
                    classroom_discussions.append(disc)
        return classroom_discussions

    # update slash command
    @bot.slash_command(
        name='upcoming',
        description="```/upcoming [end_date]``` - Updates 'Upcoming' category to show all work due before the end date")
    async def update_upcoming(ctx: discord.ApplicationContext,
                              end_date: discord.Option(str, description="End date in the format YYYY-MM-DD") = (
                                      datetime.date.today() + datetime.timedelta(days=7)).strftime('%Y-%m-%d')):

        await ctx.defer(ephemeral=True)

        upcoming_category = discord.utils.get(ctx.guild.categories, name='Upcoming')
        if upcoming_category is None:
            upcoming_category = await ctx.guild.create_category("Upcoming")

        # clears the current category so that it does not get bloated
        for channel in upcoming_category.channels:
            await channel.delete()

        # runs the get assignment data function
        assignment_date_data = await get_dates_assignments(str(ctx.guild_id), end_date)

        # runs the get date data for quizzes
        quiz_date_data = await get_dates_quiz(str(ctx.guild_id), end_date)

        discussion_date_data = await get_dates_discussion(str(ctx.guild_id), end_date)

        # iterates through all dates collected for quizzes
        for item in quiz_date_data:
            channel_name = str(item['title'])
            points = str(item['points'])

            await ctx.guild.create_voice_channel(
                name=f"❓| {points} Pts. | {channel_name}",
                category=upcoming_category
            )

        # iterates through all dates collected for assignments
        for item in assignment_date_data:
            channel_name = str(item['title'])
            points = str(item['points'])

            await ctx.guild.create_voice_channel(
                name=f"📝| {points} Pts. | {channel_name}",
                category=upcoming_category
            )

        for item in discussion_date_data:
            channel_name = str(item['title'])
            points = str(item['points'])

            await ctx.guild.create_voice_channel(
                name=f"💬| {points} Pts. | {channel_name}",
                category=upcoming_category
            )

        # adds the icon for the channels and permissions
        student_role = discord.utils.get(ctx.guild.roles, name="Student")
        assistant_role = discord.utils.get(ctx.guild.roles, name="Assistant")
        educator_role = discord.utils.get(ctx.guild.roles, name="Educator")

        await upcoming_category.set_permissions(ctx.guild.default_role, connect=False, view_channel=True)
        await upcoming_category.set_permissions(student_role, connect=False, view_channel=True)
        await upcoming_category.set_permissions(assistant_role, connect=False, view_channel=True)
        await upcoming_category.set_permissions(educator_role, connect=False, view_channel=True)
        await upcoming_category.set_permissions(ctx.guild.owner, connect=False, view_channel=True)

        await ctx.respond(f"Upcoming category updated")

    grade = bot.create_group("grade", "Grading commands for Educators")

    @grade.command(name='quiz', description='```/grade quiz [score] (comments)``` - Grades a student quiz submission')
    async def quiz(ctx: discord.ApplicationContext, score: int, comments: str = None):

        user_roles = [role.name for role in ctx.author.roles]

        await ctx.defer()

        if not any(role in ["Educator", "Assistant"] for role in user_roles):
            return await ctx.respond("You need to be an educator/assistant to grade school work", delete_after=3)

        category = ctx.channel.category.name

        if category != "Submissions":
            return await ctx.respond("This command only works in the Submissions category", delete_after=3)

        message_history = ctx.channel.history(oldest_first=True)
        first_message = await message_history.next()
        quiz_content = await message_history.next()
        get_id = first_message.content.split("ID: ")[1]
        id_values = get_id.split('-')

        if id_values[0] != 'Q':
            return await ctx.respond("This is not a quiz submission", delete_after=3)

        response = supabase.table('Quiz').select('*').eq('id', id_values[1]).execute()
        quiz_data = response.data[0]
        request = await api.get_user_id(ctx.user.id)
        grader_id = request['id']
        response = supabase.table('User').select('discordId').eq('id', id_values[2]).execute()
        student = await bot.fetch_user(response.data[0]['discordId'])
        grades_category = discord.utils.get(ctx.guild.categories, name="Grades")
        channel_name = utils.to_discord_channel_name(f"{student.display_name} {student.id % 10000}")
        grades_channel = discord.utils.get(grades_category.channels, name=channel_name)
        if grades_channel:
            message = await grades_channel.send(f"-------------------------------------------------\n"
                                        f"```You received a {score}/{quiz_data['points']} on the Quiz: {quiz_data['title']}\n\nGrader Comments: {comments}```\n\n"
                                      f"{quiz_content.content}")
        else:
            grades_channel = await ctx.guild.create_text_channel(name=channel_name, category=grades_category)
            member = ctx.guild.get_member(student.id)
            await grades_channel.set_permissions(member, view_channel=True)
            message = await grades_channel.send(f"-------------------------------------------------\n"
                f"```You received a {score}/{quiz_data['points']} on the Quiz: {quiz_data['title']}\n\nGrader Comments: {comments}```\n\n"
                                      f"{quiz_content.content}")
        data = {'studentId': id_values[2], 'taskType': 'Quiz', 'taskId': id_values[1], 'score': score,
                'graderId': grader_id, 'messageId': message.id}
        supabase.table('Grade').insert(data).execute()
        await ctx.channel.delete()

    @grade.command(name='assignment', description='```/grade assignment [score] (comments)``` - Grades a student assignment submission')
    async def assignment(ctx: discord.ApplicationContext, score: int, comments: str = None):

        user_roles = [role.name for role in ctx.author.roles]

        await ctx.defer()

        if not any(role in ["Educator", "Assistant"] for role in user_roles):
            return await ctx.respond("You need to be an educator/assistant to grade school work", delete_after=3)

        category = ctx.channel.category.name

        if category != "Submissions":
            return await ctx.respond("This command only works in the Submissions category", delete_after=3)

        message_history = ctx.channel.history(oldest_first=True)
        first_message = await message_history.next()
        assignment_content = await message_history.next()
        get_id = first_message.content.split("ID: ")[1]
        id_values = get_id.split('-')

        if id_values[0] != 'A':
            return await ctx.respond("This is not an assignment submission", delete_after=3)

        response = supabase.table('Assignment').select('*').eq('id', id_values[1]).execute()
        ass_data = response.data[0]
        request = await api.get_user_id(ctx.user.id)
        grader_id = request['id']
        response = supabase.table('User').select('discordId').eq('id', id_values[2]).execute()
        student = await bot.fetch_user(response.data[0]['discordId'])
        grades_category = discord.utils.get(ctx.guild.categories, name="Grades")
        channel_name = utils.to_discord_channel_name(f"{student.display_name} {student.id % 10000}")
        grades_channel = discord.utils.get(grades_category.channels, name=channel_name)
        file_upload = None
        if assignment_content.attachments:
            attachment = assignment_content.attachments[0]
            file = await attachment.read()
            file_obj = io.BytesIO(file)
            file_upload = discord.File(file_obj, filename=attachment.filename)
        if grades_channel:
            message = await grades_channel.send(content=f"-------------------------------------------------\n"
                                                        f"```You received a {score}/{ass_data['points']} on the Assignment: {ass_data['title']}\n\nGrader Comments: {comments}```\n\n"
                f"{assignment_content.content}", file=file_upload)
        else:
            grades_channel = await ctx.guild.create_text_channel(name=channel_name, category=grades_category)
            member = ctx.guild.get_member(student.id)
            await grades_channel.set_permissions(member, view_channel=True)
            message = await grades_channel.send(f"-------------------------------------------------\n"
                f"```You received a {score}/{ass_data['points']} on the Assignment: {ass_data['title']}\n\nGrader Comments: {comments}```\n\n"
                f"{assignment_content.content}", file=file_upload)
        data = {'studentId': id_values[2], 'taskType': 'Assignment', 'taskId': id_values[1], 'score': score,
                'graderId': grader_id, 'messageId': message.id}
        supabase.table('Grade').insert(data).execute()
        await ctx.channel.delete()

    @grade.command(name='discussion', description='```/grade assignment [score] (comments)``` - Grades a student assignment submission')
    async def discussion(ctx: discord.ApplicationContext, score: int, comments: str = None, student: discord.Member = None):
        user_roles = [role.name for role in ctx.author.roles]

        await ctx.defer(ephemeral=True)

        if not any(role in ["Educator", "Assistant"] for role in user_roles):
            return await ctx.respond("You need to be an educator to grade school work", delete_after=3)

        category = ctx.channel.category.name

        if category != "Discussions":
            return await ctx.respond("This command only works in the Discussions category", delete_after=3)

        channel_id = ctx.channel.id
        response = supabase.table("Discussion").select('*').eq('channelId', channel_id).execute()
        request = await api.get_user_id(ctx.user.id)
        grader_id = request['id']
        title = response.data[0]['title']
        points = int(response.data[0]['points'])
        taskId = int(response.data[0]['id'])
        dueDate = response.data[0]['dueDate']
        student_role = discord.utils.get(ctx.guild.roles, name="Student")
        student_score = score
        grades_category = discord.utils.get(ctx.guild.categories, name="Grades")

        if student:
            request = await api.get_user_id(student.id)
            student_id = request['id']
            response = supabase.table('Grade').select('*').match(
                {'studentId': student_id, 'taskType': 'Discussion', 'taskId': taskId}).execute()
            data = {'taskType': 'Discussion', 'graderId': grader_id, 'taskId': taskId, 'studentId': student_id,
                    'score': student_score}
            if response.data:
                supabase.table('Grade').update(data).eq('id', response.data[0]['id']).execute()
                message_id = response.data[0]['messageId']
                channel_name = utils.to_discord_channel_name(f"{student.display_name} {student.id % 10000}")
                grades_channel = discord.utils.get(grades_category.channels, name=channel_name)
                if grades_channel is None:
                    grades_channel = await ctx.guild.create_text_channel(name=channel_name, category=grades_category)
                    await grades_channel.set_permissions(student, view_channel=True)
                grade_message = await grades_channel.fetch_message(message_id)
                await grade_message.edit(content=f"-------------------------------------------------\n"
                                         f"```You received a {student_score}/{points} on the Discussion: {title}\n\nGrader Comments: {comments}```")
                return await ctx.respond(content=f"Updated Score - Gave {student.display_name} a grade of {student_score}/{points}", ephemeral=True)
            message = await student.send(f"-------------------------------------------------\n"
                f"```You received a {student_score}/{points} on the Discussion: {title}\n\nGrader Comments: {comments}```")
            data['messageId'] = message.id
            supabase.table('Grade').insert(data).execute()
            return await ctx.respond(content=f"Gave {student.display_name} a grade of {student_score}/{points}", ephemeral=True)

        # makes list of all users in channel
        users = []
        async for message in ctx.channel.history(limit=None):
            if message.author.bot:
                continue
            users.append(message.author)

        # removes duplicated users
        users = list(set(users))

        for user in users:
            if student_role in user.roles:
                request = await api.get_user_id(user.id)
                student_id = request['id']
                data = {'taskType': 'Discussion', 'graderId': grader_id, 'taskId': taskId, 'studentId': student_id, 'score': student_score}
                response = supabase.table('Grade').select('*').match(
                    {'studentId': student_id, 'taskType': 'Discussion', 'taskId': taskId}).execute()
                channel_name = utils.to_discord_channel_name(f"{user.display_name} {user.id % 10000}")
                grades_channel = discord.utils.get(grades_category.channels, name=channel_name)
                if grades_channel is None:
                    grades_channel = await ctx.guild.create_text_channel(name=channel_name, category=grades_category)
                    await grades_channel.set_permissions(user, view_channel=True)
                if response.data:
                    supabase.table('Grade').update(data).eq('id', response.data[0]['id']).execute()
                    message_id = response.data[0]['messageId']
                    grade_message = await grades_channel.fetch_message(message_id)
                    await grade_message.edit(content=f"-------------------------------------------------\n"
                                                     f"```You received a {student_score}/{points} on the Discussion: {title}\n\nGrader Comments: {comments}```")
                else:
                    message = await grades_channel.send(f"-------------------------------------------------\n"
                                                        f"```You received a {student_score}/{points} on the Discussion: {title}\n\nGrader Comments: {comments}```")
                    data['messageId'] = message.id
                    supabase.table('Grade').insert(data).execute()

        await ctx.respond("Graded this discussion", ephemeral=True)

        channel = ctx.channel
        due = datetime.datetime.strptime(dueDate, "%Y-%m-%d").date()
        today = datetime.datetime.today().date()

        if due < today:
            await channel.set_permissions(student_role, send_messages=False)

    @grade.command(name='edit', description='```/grade edit [link] [score] (comments)``` - Used to edit a grade')
    async def edit(ctx: discord.ApplicationContext, link: discord.Option(str, description="Copy Message Link"), score: int, comments: str = None):

        await ctx.defer(ephemeral=True)

        link_split = link.split('/')
        channel_id = link_split[-2]
        message_id = link_split[-1]

        user_roles = [role.name for role in ctx.author.roles]

        if not any(role in ["Educator", "Assistant"] for role in user_roles):
            return await ctx.respond("You need to be an educator to grade school work", delete_after=3)

        category = ctx.channel.category.name

        if category != "Grades":
            return await ctx.respond("You can only edit grades in the Grades category", delete_after=3)

        channel = ctx.channel

        if int(channel_id) != int(channel.id):
            print(channel_id, channel.id)
            return await ctx.respond("That grade link does not belong to this student", delete_after=3)

        try:
            grade_message = await channel.fetch_message(message_id)
        except discord.NotFound:
            return await ctx.respond("Could not find message from link", delete_after=3)

        grade_content = grade_message.content

        response = supabase.table('Grade').select('*').eq('messageId', message_id).execute()
        grade_obj = response.data[0]
        table = grade_obj['taskType']
        table_id = grade_obj['taskId']
        response = supabase.table(table).select('*').eq('id', table_id).execute()
        task_obj = response.data[0]
        title = task_obj['title']
        points = task_obj['points']

        edited_message = grade_content + f"\n```UPDATED GRADE: {score}/{points} for {table}: {title}\n\nGrader Comments: {comments}```"

        grade_obj['score'] = score

        supabase.table('Grade').update(grade_obj).eq('messageId', message_id).execute()

        await grade_message.edit(content=edited_message)

        await ctx.respond("Grade Updated", delete_after=3)

    @bot.slash_command(name='submit',
                       description='```/submit assignment (file) (url)``` - student submit assignment')
    async def submit(ctx: discord.ApplicationContext, file: discord.Attachment = None, url: str = None):

        await ctx.defer(ephemeral=True)

        student_role = discord.utils.get(ctx.guild.roles, name="Student")
        if student_role not in ctx.author.roles:
            return await ctx.respond("Only Students can use /submit command")

        category_name = ctx.channel.category.name
        if category_name != "Assignments":
            return await ctx.respond("/submit is only allowed in the designated assignment channel.")

        if file and url:
            return await ctx.respond("Please provide either a file or a URL, not both.")

        if not file and not url:
            return await ctx.respond(
                "Please provide either a file or a URL.\nIf the file you are trying to upload is too large is the `/upload_file` command")

        first_message = await ctx.channel.history(oldest_first=True, limit=1).next()
        embed = first_message.embeds[0]
        due_date_str = embed.fields[4].value
        due_date = datetime.datetime.strptime(due_date_str, "%Y-%m-%d").date()
        today = datetime.date.today()

        discord_id = ctx.author.id
        studentId_dict = await api.get_user_id(discord_id)
        studentId = studentId_dict["id"]

        chan_id = ctx.channel_id
        assignment_dict = await api.get_assignment(chan_id)
        assignmentId = assignment_dict['Assignment'].id
        assignmentTitle = assignment_dict['Assignment'].title

        # Check if a Submission category already exists
        submission_category = discord.utils.get(ctx.guild.categories, name="Submissions")

        if submission_category is None:
            # Create the Submission category if it doesn't exist
            submission_category = await ctx.guild.create_category("Submissions")

        # Create the text channel with the name "📝assignmenttitle-studentname" under the Submission category
        channel_name = f"📝{assignment_dict['Assignment'].title}-{ctx.author.name}"
        # channel = await ctx.guild.create_text_channel(channel_name, category=submission_category)
        existing_channel = discord.utils.get(submission_category.text_channels, name=channel_name)
        for channel in submission_category.channels:
            if channel.name == utils.to_discord_channel_name(channel_name):
                await channel.delete()

        channel = await ctx.guild.create_text_channel(channel_name, category=submission_category)

        await channel.send(f"ID: A-{assignmentId}-{studentId}")
        if today > due_date:
            await channel.send(f"```diff\n- THIS SUBMISSION IS {abs((today - due_date).days)} DAY(S) LATE```")
        if file:
            await channel.send(
                content=f"**Assignment - {assignmentTitle}**\t{assignment_dict['Assignment'].points} Pts.\nStudent: {ctx.author.display_name}\n",
                file=await file.to_file())
        else:
            await channel.send(
                content=f"**Assignment - {assignmentTitle}**\t{assignment_dict['Assignment'].points} Pts.\nStudent: {ctx.author.display_name}\n\nSubmission: {url}")

        return await ctx.respond("Assignment Submitted!")

    bot.run(DISCORD_TOKEN)
